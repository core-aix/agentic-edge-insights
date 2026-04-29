#!/usr/bin/env python3
import json
import math
import os
import ast
import tempfile
import io
import logging
import warnings
import contextlib
import shutil
import sys
import zipfile
from collections import Counter, defaultdict, deque
from pathlib import Path

import pandas as pd


def _resolve_path(path: str | None, env_dir_var: str) -> Path | None:
    p = str(path or "").strip()
    if p:
        q = Path(p)
        if q.exists():
            return q
    d = str(os.environ.get(env_dir_var, "") or "").strip()
    if d and Path(d).exists():
        return Path(d)
    return None


def _read_table(path: Path) -> pd.DataFrame:
    if path.suffix.lower() == ".tsv":
        return pd.read_csv(path, sep="\t")
    return pd.read_csv(path)


def _find_first(base: Path, names: list[str]) -> Path | None:
    if base.is_file():
        return base
    for n in names:
        p = base / n
        if p.exists() and p.is_file():
            return p
    return None


def _find_topology_path(base: Path) -> Path | None:
    if base.is_file() and base.suffix.lower() == ".json":
        return base
    if not base.exists() or not base.is_dir():
        return None
    for p in base.rglob("*.json"):
        n = p.name.lower()
        if "topology" in n or "graph" in n:
            return p
    return None


def _safe_json_or_literal(text: str):
    s = str(text or "").strip()
    if not s:
        return None
    try:
        return json.loads(s)
    except Exception:
        pass
    try:
        return ast.literal_eval(s)
    except Exception:
        return None


def _extract_candidate_entities_from_sre(base: Path) -> dict:
    events_path = _find_first(base, ["k8s_events_raw.tsv"])
    objects_path = _find_first(base, ["k8s_objects_raw.tsv"])
    id_counter: Counter[str] = Counter()
    name_counter: Counter[str] = Counter()
    namespace_counter: Counter[str] = Counter()

    def add_if(v: object, ctr: Counter[str]):
        s = str(v or "").strip()
        if s:
            ctr[s] += 1

    if events_path and events_path.exists():
        try:
            ev = pd.read_csv(events_path, sep="\t")
            if "Body" in ev.columns:
                for raw in ev["Body"].astype(str).head(2000):
                    obj = _safe_json_or_literal(raw)
                    if not isinstance(obj, dict):
                        continue
                    o = obj.get("object", obj)
                    if not isinstance(o, dict):
                        continue
                    md = (
                        o.get("metadata", {})
                        if isinstance(o.get("metadata", {}), dict)
                        else {}
                    )
                    regarding = (
                        o.get("regarding", {})
                        if isinstance(o.get("regarding", {}), dict)
                        else {}
                    )
                    involved = (
                        o.get("involvedObject", {})
                        if isinstance(o.get("involvedObject", {}), dict)
                        else {}
                    )
                    for d in [md, regarding, involved]:
                        add_if(d.get("uid"), id_counter)
                        add_if(d.get("name"), name_counter)
                        add_if(d.get("namespace"), namespace_counter)
            if "LogAttributes" in ev.columns:
                for raw in ev["LogAttributes"].astype(str).head(2000):
                    la = _safe_json_or_literal(raw)
                    if isinstance(la, dict):
                        add_if(la.get("k8s.namespace.name"), namespace_counter)
        except Exception:
            pass

    if objects_path and objects_path.exists():
        try:
            ob = pd.read_csv(objects_path, sep="\t")
            if "Body" in ob.columns:
                for raw in ob["Body"].astype(str).head(2000):
                    obj = _safe_json_or_literal(raw)
                    if not isinstance(obj, dict):
                        continue
                    o = obj.get("object", obj)
                    if not isinstance(o, dict):
                        continue
                    md = (
                        o.get("metadata", {})
                        if isinstance(o.get("metadata", {}), dict)
                        else {}
                    )
                    add_if(md.get("uid"), id_counter)
                    add_if(md.get("name"), name_counter)
                    add_if(md.get("namespace"), namespace_counter)
        except Exception:
            pass

    return {
        "events_file": str(events_path) if events_path else "",
        "objects_file": str(objects_path) if objects_path else "",
        "candidate_entity_ids": [k for k, _ in id_counter.most_common(80)],
        "candidate_entity_names": [k for k, _ in name_counter.most_common(120)],
        "candidate_namespaces": [k for k, _ in namespace_counter.most_common(40)],
    }


def _json(obj) -> str:
    return json.dumps(obj, ensure_ascii=False)


# ITBench tools
def nl2kubectl(nl_query: str) -> str:
    base = _resolve_path("", "ITBENCH_SCENARIO_DIR")
    q = str(nl_query or "").lower()
    if not base:
        return "No ITBench scenario directory set (ITBENCH_SCENARIO_DIR)."
    m = None
    q_raw = str(nl_query or "")
    m = __import__("re").search(
        r"inspect file\s+(.+)$", q_raw, flags=__import__("re").IGNORECASE
    )
    if m:
        p = Path(m.group(1).strip())
        if p.exists() and p.is_file():
            try:
                return _read_table(p).head(40).to_csv(index=False)
            except Exception:
                return p.read_text(encoding="utf-8", errors="ignore")[:12000]
    obj = _find_first(base, ["k8s_objects_raw.tsv"])
    ev = _find_first(base, ["k8s_events_raw.tsv"])
    if "pod" in q and obj is not None:
        df = _read_table(obj)
        cols = {c.lower(): c for c in df.columns}
        out = df
        if "kind" in cols:
            out = out[out[cols["kind"]].astype(str).str.lower() == "pod"]
        return out.head(40).to_csv(index=False)
    if "event" in q and ev is not None:
        return _read_table(ev).head(40).to_csv(index=False)
    if obj is not None:
        summary = _extract_candidate_entities_from_sre(base)
        return _json(summary)
    return _json({"error": "No kubernetes snapshot files found"})


def query_loki_logs(query: str, limit: int = 100, **kwargs) -> str:
    base = _resolve_path("", "ITBENCH_SCENARIO_DIR")
    if not base:
        return "No ITBench scenario directory set (ITBENCH_SCENARIO_DIR)."
    terms = [
        t for t in str(query or "").replace("{", " ").replace("}", " ").split() if t
    ]
    lines = []
    for p in base.rglob("*.log"):
        try:
            for line in p.read_text(encoding="utf-8", errors="ignore").splitlines():
                s = line.lower()
                if all(t.lower() in s for t in terms[:4]):
                    lines.append({"file": str(p), "line": line})
                if len(lines) >= max(1, int(limit)):
                    return _json(lines)
        except Exception:
            continue
    return _json(lines)


def query_jaeger_traces(
    service: str, operation: str = "", limit: int = 20, **kwargs
) -> str:
    base = _resolve_path("", "ITBENCH_SCENARIO_DIR")
    if not base:
        return "No ITBench scenario directory set (ITBENCH_SCENARIO_DIR)."
    out = []
    for p in base.rglob("*trace*.json"):
        try:
            data = json.loads(p.read_text(encoding="utf-8", errors="ignore"))
        except Exception:
            continue
        rows = (
            data
            if isinstance(data, list)
            else data.get("data", [])
            if isinstance(data, dict)
            else []
        )
        for r in rows:
            text = json.dumps(r, ensure_ascii=False).lower()
            if str(service or "").lower() in text and (
                not operation or str(operation).lower() in text
            ):
                out.append(r)
            if len(out) >= max(1, int(limit)):
                return _json(out)
    return _json(out)


def get_alerts(**kwargs) -> str:
    base = _resolve_path("", "ITBENCH_SCENARIO_DIR")
    if not base:
        return "No ITBench scenario directory set (ITBENCH_SCENARIO_DIR)."
    out = []
    for p in base.rglob("*alert*"):
        if p.is_file() and p.suffix.lower() in {
            ".json",
            ".yaml",
            ".yml",
            ".txt",
            ".log",
        }:
            out.append(str(p))
    return _json({"alerts_files": out[:100]})


def _load_topology(topology: str) -> tuple[dict, list[tuple[str, str]]]:
    p = Path(topology)
    data = json.loads(p.read_text(encoding="utf-8", errors="ignore"))
    nodes = {}
    edges = []
    raw_nodes = (
        data.get("nodes", data.get("vertices", [])) if isinstance(data, dict) else []
    )
    raw_edges = (
        data.get("edges", data.get("links", [])) if isinstance(data, dict) else []
    )
    for n in raw_nodes:
        if isinstance(n, dict):
            nid = str(n.get("id", n.get("node_id", n.get("name", ""))))
            if nid:
                nodes[nid] = n
    for e in raw_edges:
        if isinstance(e, dict):
            s = str(e.get("source", e.get("from", "")))
            t = str(e.get("target", e.get("to", "")))
            if s and t:
                edges.append((s, t))
                edges.append((t, s))
    return nodes, edges


def get_topology_nodes(topology: str = "") -> str:
    p = _resolve_path(topology, "ITBENCH_TOPOLOGY_PATH")
    if not p:
        base = _resolve_path("", "ITBENCH_SCENARIO_DIR")
        if base:
            p = _find_topology_path(base)
    if not p:
        return "No topology path provided"
    nodes, _ = _load_topology(str(p))
    return _json({"count": len(nodes), "nodes": list(nodes.keys())[:500]})


def walk_path(
    topology: str, start_id: str, start_node_type: str, target_node_type: str
) -> str:
    nodes, edges = _load_topology(topology)
    g = defaultdict(list)
    for a, b in edges:
        g[a].append(b)
    q = deque([str(start_id)])
    seen = {str(start_id)}
    hits = []
    while q and len(hits) < 50:
        u = q.popleft()
        n = nodes.get(u, {})
        t = str(n.get("type", n.get("kind", ""))).lower()
        if str(target_node_type).lower() == t and u != str(start_id):
            hits.append(u)
        for v in g.get(u, []):
            if v not in seen:
                seen.add(v)
                q.append(v)
    return _json({"targets": hits})


def get_node_info_by_name(topology: str, node_name: str) -> str:
    nodes, _ = _load_topology(topology)
    name = str(node_name or "").lower()
    for _, n in nodes.items():
        if name == str(n.get("name", "")).lower():
            return _json(n)
    return _json({})


def get_neighbors(topology: str, node_id: str) -> str:
    _, edges = _load_topology(topology)
    out = [b for a, b in edges if a == str(node_id)]
    return _json({"neighbors": out[:200]})


def check_directly_connected(topology: str, node_id1: str, node_id2: str) -> str:
    _, edges = _load_topology(topology)
    s = {(a, b) for a, b in edges}
    return _json({"connected": (str(node_id1), str(node_id2)) in s})


def summarize_sre_candidates(scenario_dir: str = "") -> str:
    base = _resolve_path(scenario_dir, "ITBENCH_SCENARIO_DIR")
    if not base:
        return _json(
            {"error": "No ITBench scenario directory set (ITBENCH_SCENARIO_DIR)."}
        )
    summary = _extract_candidate_entities_from_sre(base)
    return _json(summary)


def rank_sre_root_cause_candidates(scenario_dir: str = "") -> str:
    base = _resolve_path(scenario_dir, "ITBENCH_SCENARIO_DIR")
    if not base:
        return _json(
            {"error": "No ITBench scenario directory set (ITBENCH_SCENARIO_DIR)."}
        )

    summary = _extract_candidate_entities_from_sre(base)
    ids = summary.get("candidate_entity_ids", []) or []
    names = summary.get("candidate_entity_names", []) or []

    events_path = _find_first(base, ["k8s_events_raw.tsv"])
    score = Counter()
    uuid_re = __import__("re").compile(
        r"^[0-9a-f]{8}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{12}$",
        flags=__import__("re").IGNORECASE,
    )

    def _is_uuid_like(v: str) -> bool:
        return bool(uuid_re.match(str(v or "").strip()))

    if events_path and events_path.exists():
        try:
            ev = pd.read_csv(events_path, sep="\t")
            for raw in ev.get("Body", pd.Series([], dtype=str)).astype(str).head(3000):
                txt = raw.lower()
                weight = 1
                if any(
                    k in txt for k in ["error", "fail", "panic", "crash", "backoff"]
                ):
                    weight = 4
                elif any(
                    k in txt for k in ["warning", "unhealthy", "denied", "forbidden"]
                ):
                    weight = 2
                if weight <= 1:
                    continue
                for i in ids[:120]:
                    if str(i).lower() in txt:
                        score[i] += weight
                for n in names[:160]:
                    nl = str(n).lower()
                    if nl and nl in txt:
                        score[n] += 2 * weight
        except Exception:
            pass

    ranked_all = [k for k, _ in score.most_common(80)]
    if not ranked_all:
        ranked_all = names[:40] + ids[:40]

    # Filter likely observability/control-plane noise values.
    deny_exact = {
        "default",
        "kube-system",
        "ingress-nginx",
        "monitoring",
        "logging",
        "observability",
        "ad",
    }
    deny_substrings = [
        "topology-monitor",
        "otel-collector",
        "prometheus-alert",
        "node-exporter",
        "grafana",
        "jaeger",
        "loki",
        "root-cause_entity_group_id",
    ]

    def _is_noise(v: str) -> bool:
        s = str(v or "").strip().lower()
        if not s:
            return True
        if s in deny_exact:
            return True
        return any(tok in s for tok in deny_substrings)

    def _is_low_value(v: str) -> bool:
        s = str(v or "").strip().lower()
        if _is_noise(s):
            return True
        if s.startswith("pvc-"):
            return True
        if _is_uuid_like(s):
            return True
        return False

    ranked_names = [x for x in ranked_all if not _is_low_value(str(x))]
    # Fall back to id-like entities only if we have no useful names.
    if not ranked_names:
        ranked_names = [x for x in ranked_all if not _is_noise(str(x))]

    return _json(
        {
            "ranked_candidates": ranked_names[:20],
            "candidate_entity_ids": ids[:30],
            "candidate_entity_names": names[:50],
        }
    )


def analyze_finops_cost_anomaly(
    path: str, anomaly_date: str = "", account_id: str = "", top_k: int = 5
) -> str:
    def parse_mixed_dates(values):
        try:
            return pd.to_datetime(values, errors="coerce", format="mixed")
        except Exception:
            return pd.to_datetime(values, errors="coerce")

    p = Path(str(path or "").strip())
    if not p.exists() or p.is_dir():
        return _json({"error": f"File not found: {path}"})

    try:
        df = _read_table(p)
    except Exception as exc:
        return _json({"error": f"Failed to read table: {exc}"})

    cols = {str(c).strip().lower(): c for c in df.columns}

    def pick_col(cands: list[str]) -> str:
        for c in cands:
            if c in cols:
                return cols[c]
        return ""

    date_col = pick_col(["date", "timestamp", "event_date"])
    account_col = pick_col(["account_id", "account id", "account"])
    resource_col = pick_col(
        [
            "instance_family",
            "instance_family_name",
            "resource_name",
            "resource",
            "instance_id",
            "instance",
        ]
    )
    cost_col = pick_col(["unblended_cost", "cost", "total_cost", "amount"])

    if not resource_col or not cost_col:
        return _json(
            {
                "error": "Missing required columns",
                "required": ["resource column", "cost column"],
                "detected_columns": [str(c) for c in df.columns],
            }
        )

    work = df.copy()
    work[resource_col] = work[resource_col].astype(str).str.strip()
    work[cost_col] = pd.to_numeric(work[cost_col], errors="coerce").fillna(0.0)

    parsed_anomaly = None
    ad = str(anomaly_date or "").strip()
    if ad:
        try:
            parsed_anomaly = parse_mixed_dates(ad)
            if pd.isna(parsed_anomaly):
                parsed_anomaly = None
        except Exception:
            parsed_anomaly = None

    date_match_found = False
    available_dates = []
    if date_col:
        work["__normalized_date"] = parse_mixed_dates(work[date_col])
        valid_dates = work["__normalized_date"].dropna()
        if len(valid_dates) > 0:
            available_dates = sorted(
                {str(d.date()) for d in valid_dates.head(2000).tolist()}
            )[:20]
        if parsed_anomaly is not None:
            d = parsed_anomaly.date()
            mask = work["__normalized_date"].dt.date == d
            date_match_found = bool(mask.any())
            if date_match_found:
                work = work[mask]

    used_account_filter = False
    acct = str(account_id or "").strip()
    if acct and account_col:
        m = work[account_col].astype(str).str.strip() == acct
        if m.any():
            work = work[m]
            used_account_filter = True

    if len(work) == 0:
        return _json(
            {
                "error": "No matching rows after filters",
                "date_match_found": date_match_found,
                "used_account_filter": used_account_filter,
                "normalized_anomaly_date": str(parsed_anomaly.date())
                if parsed_anomaly is not None
                else "",
                "available_dates": available_dates,
            }
        )

    grouped = (
        work.groupby(resource_col, dropna=False)[cost_col]
        .sum()
        .sort_values(ascending=False)
        .reset_index()
    )
    grouped.columns = ["resource", "total_cost"]

    k = max(1, int(top_k or 5))
    top_rows = grouped.head(k)
    total_cost = float(grouped["total_cost"].sum())
    top_resources = top_rows["resource"].astype(str).tolist()
    top_costs = [float(x) for x in top_rows["total_cost"].tolist()]

    mean_cost = float(grouped["total_cost"].mean()) if len(grouped) else 0.0
    std_cost = float(grouped["total_cost"].std(ddof=0)) if len(grouped) > 1 else 0.0
    if std_cost > 0:
        strong = grouped[grouped["total_cost"] >= (mean_cost + std_cost)]
        strong_candidates = strong["resource"].astype(str).head(k).tolist()
    else:
        strong_candidates = []
    if not strong_candidates:
        strong_candidates = top_resources[: max(1, min(3, len(top_resources)))]

    return _json(
        {
            "normalized_anomaly_date": str(parsed_anomaly.date())
            if parsed_anomaly is not None
            else "",
            "date_match_found": date_match_found,
            "available_dates": available_dates,
            "resource_column": resource_col,
            "cost_column": cost_col,
            "rows_used": int(len(work)),
            "top_resources_by_cost": top_resources,
            "top_costs": top_costs,
            "strong_candidate_resources": strong_candidates,
            "recommended_answer": ", ".join(strong_candidates),
            "total_cost_considered": total_cost,
        }
    )


TOOL_REGISTRY = {
    "analyze_finops_cost_anomaly": analyze_finops_cost_anomaly,
    "nl2kubectl": nl2kubectl,
    "query_loki_logs": query_loki_logs,
    "query_jaeger_traces": query_jaeger_traces,
    "get_alerts": get_alerts,
    "get_topology_nodes": get_topology_nodes,
    "walk_path": walk_path,
    "get_node_info_by_name": get_node_info_by_name,
    "get_neighbors": get_neighbors,
    "check_directly_connected": check_directly_connected,
    "summarize_sre_candidates": summarize_sre_candidates,
    "rank_sre_root_cause_candidates": rank_sre_root_cause_candidates,
}


def run_tool(tool: str, args: dict) -> str:
    fn = TOOL_REGISTRY.get(str(tool or ""))
    if not fn:
        return _json({"error": f"Unknown benchmark tool: {tool}"})
    try:
        return fn(**(args or {}))
    except TypeError as exc:
        return _json({"error": f"Invalid args for {tool}: {exc}"})
    except Exception as exc:  # noqa: BLE001
        return _json({"error": f"{tool} execution error: {exc}"})


DOCLING_QUIET_OVERRIDE: bool | None = None
_DOCLING_CONVERTER = None


def configure_docling_quiet(enabled: bool | None = None) -> None:
    global DOCLING_QUIET_OVERRIDE
    DOCLING_QUIET_OVERRIDE = enabled


def _docling_quiet_enabled() -> bool:
    if DOCLING_QUIET_OVERRIDE is not None:
        return bool(DOCLING_QUIET_OVERRIDE)
    return str(os.environ.get("BENCHMARK_DOCLING_QUIET", "1")).strip().lower() not in {
        "0",
        "false",
        "no",
        "off",
    }


def _select_docling_ocr_options():
    from docling.datamodel.pipeline_options import (
        OcrAutoOptions,
        OcrMacOptions,
        TesseractCliOcrOptions,
    )

    if sys.platform == "darwin":
        return OcrMacOptions()
    if shutil.which("tesseract"):
        return TesseractCliOcrOptions()
    return OcrAutoOptions()


def _get_docling_converter():
    global _DOCLING_CONVERTER
    if _DOCLING_CONVERTER is not None:
        return _DOCLING_CONVERTER

    from docling.document_converter import (
        DocumentConverter,
        ImageFormatOption,
        PdfFormatOption,
    )
    from docling.datamodel.base_models import InputFormat
    from docling.datamodel.pipeline_options import PdfPipelineOptions

    ocr_opts = _select_docling_ocr_options()
    pdf_opts = PdfPipelineOptions(ocr_options=ocr_opts)
    fmt_opts = {
        InputFormat.PDF: PdfFormatOption(pipeline_options=pdf_opts),
        InputFormat.IMAGE: ImageFormatOption(pipeline_options=pdf_opts),
    }
    _DOCLING_CONVERTER = DocumentConverter(format_options=fmt_opts)
    return _DOCLING_CONVERTER


def _truncate(text: str, max_chars: int = 6000) -> str:
    text = str(text)
    if len(text) <= max_chars:
        return text
    return text[:max_chars] + "\n...[truncated]"


def _docling_extract_text(path: str, max_chars: int = 12000) -> str:
    try:
        converter = _get_docling_converter()
    except Exception as exc:  # noqa: BLE001
        return f"Read file error: docling unavailable: {exc}"

    try:
        quiet = _docling_quiet_enabled()
        noisy_loggers = [
            "docling",
            "docling_core",
            "docling_parse",
            "rapidocr",
            "ocrmac",
        ]
        old_levels: dict[str, int] = {}
        if quiet:
            for name in noisy_loggers:
                lg = logging.getLogger(name)
                old_levels[name] = lg.level
                lg.setLevel(logging.ERROR)

        try:
            if quiet:
                with warnings.catch_warnings():
                    warnings.filterwarnings(
                        "ignore",
                        category=DeprecationWarning,
                        module=r"docling\..*",
                    )
                    warnings.filterwarnings(
                        "ignore",
                        message=r".*Clashing hyperlinks.*",
                        category=Warning,
                    )
                    warnings.filterwarnings(
                        "ignore",
                        message=r".*Clashing formatting.*",
                        category=Warning,
                    )
                    with contextlib.redirect_stderr(io.StringIO()):
                        with contextlib.redirect_stdout(io.StringIO()):
                            result = converter.convert(str(path))
            else:
                result = converter.convert(str(path))
        finally:
            if quiet:
                for name, lvl in old_levels.items():
                    logging.getLogger(name).setLevel(lvl)

        doc = result.document
        txt = ""
        export_md = getattr(doc, "export_to_markdown", None)
        export_txt = getattr(doc, "export_to_text", None)
        if callable(export_md):
            txt = str(export_md() or "").strip()
        if not txt and callable(export_txt):
            txt = str(export_txt() or "").strip()
        if not txt:
            txt = str(doc)
        return _truncate(txt, max_chars)
    except Exception as exc:  # noqa: BLE001
        return f"Read file error: docling extraction failed: {exc}"


def _read_zip_archive(path: str, max_chars: int = 12000) -> str:
    p = Path(path)
    if not p.exists() or p.is_dir():
        return f"Read file error: invalid zip path: {path}"
    try:
        with tempfile.TemporaryDirectory() as td:
            extract_root = Path(td) / "unzipped"
            extract_root.mkdir(parents=True, exist_ok=True)
            with zipfile.ZipFile(str(p), "r") as zf:
                zf.extractall(str(extract_root))

            files = [x for x in extract_root.rglob("*") if x.is_file()]
            if not files:
                return "ZIP archive is empty"

            blocks = [f"ZIP entries: {len(files)}"]
            budget = max_chars
            for fp in sorted(files)[:20]:
                rel = str(fp.relative_to(extract_root))
                try:
                    txt = read_file(str(fp), max_chars=min(3000, max(400, budget // 4)))
                except Exception as exc:  # noqa: BLE001
                    txt = f"Read file error: {exc}"
                blk = f"=== FILE: {rel} ===\n{txt}"
                blocks.append(blk)
                budget -= len(blk)
                if budget <= 0:
                    break
            return _truncate("\n\n".join(blocks), max_chars)
    except Exception as exc:  # noqa: BLE001
        return f"Read file error: zip extraction failed: {exc}"


def read_file(path: str, max_chars: int = 12000) -> str:
    p = Path(path)
    if not p.exists():
        return f"File not found: {path}"
    if p.is_dir():
        return f"Path is a directory, not a file: {path}"
    suffix = p.suffix.lower()
    try:
        if suffix in {".txt", ".md", ".json", ".yaml", ".yml", ".log", ".csv", ".tsv"}:
            return _truncate(p.read_text(encoding="utf-8", errors="ignore"), max_chars)
        if suffix in {".xls", ".xlsx"}:
            from openpyxl import load_workbook

            out = []
            wb = load_workbook(str(p), data_only=True)
            for ws in wb.worksheets[:5]:
                rows = []
                for row in ws.iter_rows(min_row=1, max_row=80, max_col=30, values_only=True):
                    vals = ["" if v is None else str(v) for v in row]
                    if any(vals):
                        rows.append(",".join(vals))
                out.append(f"=== SHEET: {ws.title} ({ws.max_row}x{ws.max_column}) ===\n" + "\n".join(rows[:120]))
            return _truncate("\n\n".join(out), max_chars)
        if suffix == ".pdf":
            return _docling_extract_text(str(p), max_chars=max_chars)
        if suffix == ".docx":
            import docx

            d = docx.Document(str(p))
            txt = "\n".join(par.text for par in d.paragraphs)
            return _truncate(txt, max_chars)
        if suffix == ".pptx":
            from pptx import Presentation

            prs = Presentation(str(p))
            blocks = []
            for i, slide in enumerate(prs.slides, start=1):
                lines = []
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text:
                        lines.append(text)
                blocks.append(f"=== SLIDE {i} ===\n" + "\n".join(lines))
            return _truncate("\n\n".join(blocks), max_chars)
        if suffix == ".zip":
            return _read_zip_archive(str(p), max_chars=max_chars)
        if suffix in {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".mp3", ".wav", ".m4a", ".flac", ".ogg"}:
            return _docling_extract_text(str(p), max_chars=max_chars)
    except Exception as exc:  # noqa: BLE001
        return f"Read file error: {exc}"
    return f"Unsupported file type: {suffix}"


def summarize_csv(
    path: str,
    group_by: str = "",
    value_column: str = "",
    top_k: int = 10,
) -> str:
    p = Path(path)
    if not p.exists():
        return f"File not found: {path}"
    if p.is_dir():
        return f"Path is a directory, not a file: {path}"
    suffix = p.suffix.lower()
    if suffix not in {".csv", ".tsv"}:
        return f"Unsupported file type for summarize_csv: {suffix}"
    try:
        sep = "\t" if suffix == ".tsv" else ","
        df = pd.read_csv(p, sep=sep)
        out = {
            "rows": int(len(df)),
            "columns": [str(c) for c in df.columns],
            "preview": df.head(5).to_dict(orient="records"),
        }
        gb = str(group_by or "").strip()
        vc = str(value_column or "").strip()
        k = max(1, min(int(top_k or 10), 100))
        if gb and gb in df.columns:
            if vc and vc in df.columns and pd.api.types.is_numeric_dtype(df[vc]):
                agg_df = (
                    df.groupby(gb, dropna=False)[vc]
                    .agg(["count", "sum", "mean", "max"])
                    .reset_index()
                )
                agg_rows = agg_df.to_dict(orient="records")
                agg_rows.sort(key=lambda r: float(r.get("sum", 0) or 0), reverse=True)
                out["grouped"] = agg_rows[:k]
            else:
                counts = df[gb].astype(str).value_counts(dropna=False).head(k).reset_index()
                counts.columns = [gb, "count"]
                out["grouped"] = counts.to_dict(orient="records")
        return _truncate(json.dumps(out, ensure_ascii=False, indent=2), 12000)
    except Exception as exc:  # noqa: BLE001
        return f"CSV summary error: {exc}"


TOOL_SCHEMAS = {
    "read_file": {"path": "local file path"},
    "summarize_csv": {
        "path": "local CSV file path",
        "group_by": "column name (optional)",
        "value_column": "numeric column name (optional)",
        "top_k": "int (optional, default=10)",
    },
    "nl2kubectl": {"nl_query": "natural-language kubectl request"},
    "query_loki_logs": {
        "query": "LogQL query",
        "limit": "int (optional)",
        "start": "string (optional)",
        "end": "string (optional)",
        "since": "string (optional)",
        "step": "string (optional)",
        "interval": "string (optional)",
        "direction": "forward|backward (optional)",
    },
    "query_jaeger_traces": {
        "service": "service name",
        "operation": "operation filter (optional)",
        "start_time": "int microseconds",
        "end_time": "int microseconds",
        "limit": "int",
        "error_traces_only": "bool (optional)",
    },
    "get_alerts": {},
    "get_topology_nodes": {},
    "walk_path": {
        "topology": "topology JSON path",
        "start_id": "start node id",
        "start_node_type": "start node type",
        "target_node_type": "target node type",
    },
    "get_node_info_by_name": {
        "topology": "topology JSON path",
        "node_name": "node name",
    },
    "get_neighbors": {
        "topology": "topology JSON path",
        "node_id": "node id",
    },
    "check_directly_connected": {
        "topology": "topology JSON path",
        "node_id1": "node id 1",
        "node_id2": "node id 2",
    },
    "analyze_finops_cost_anomaly": {
        "path": "local CSV file path",
        "anomaly_date": "date string from anomaly details (optional)",
        "account_id": "account id from anomaly details (optional)",
        "top_k": "int (optional, default=5)",
    },
    "summarize_sre_candidates": {"scenario_dir": "optional local scenario directory"},
    "rank_sre_root_cause_candidates": {
        "scenario_dir": "optional local scenario directory"
    },
}


TOOL_DESCRIPTIONS = {
    "read_file": "Reads local files (txt/csv/xlsx/pdf/docx/images/audio).",
    "summarize_csv": "Summarizes CSV columns and optional grouped numeric stats.",
    "nl2kubectl": "ITBench SRE NL-to-kubectl tool.",
    "query_loki_logs": "ITBench observability log query tool (Loki).",
    "query_jaeger_traces": "ITBench observability trace query tool (Jaeger).",
    "get_alerts": "ITBench observability alert retrieval tool.",
    "get_topology_nodes": "ITBench topology node retrieval tool.",
    "walk_path": "ITBench graph traversal path walk tool.",
    "get_node_info_by_name": "ITBench graph node info tool.",
    "get_neighbors": "ITBench graph neighbor tool.",
    "check_directly_connected": "ITBench graph connectivity check tool.",
    "analyze_finops_cost_anomaly": "ITBench FinOps anomaly helper with date normalization.",
    "summarize_sre_candidates": "ITBench SRE candidate entity summarization tool.",
    "rank_sre_root_cause_candidates": "ITBench SRE heuristic ranking of root-cause candidates.",
}


BENCHMARK_TOOLSETS = {
    "itbench_lite_finops": {
        "analyze_finops_cost_anomaly",
        "read_file",
        "summarize_csv",
    },
    "itbench_lite_sre": {
        "read_file",
        "nl2kubectl",
        "query_loki_logs",
        "query_jaeger_traces",
        "get_alerts",
        "get_topology_nodes",
        "walk_path",
        "get_node_info_by_name",
        "get_neighbors",
        "check_directly_connected",
        "summarize_sre_candidates",
        "rank_sre_root_cause_candidates",
    },
    "itbench_lite": {
        "read_file",
        "summarize_csv",
        "nl2kubectl",
        "query_loki_logs",
        "query_jaeger_traces",
        "get_alerts",
        "get_topology_nodes",
        "walk_path",
        "get_node_info_by_name",
        "get_neighbors",
        "check_directly_connected",
        "summarize_sre_candidates",
        "rank_sre_root_cause_candidates",
    },
}


BENCHMARK_TOOL_GUIDANCE = {
    "itbench_lite_finops": [
        "This is an ITBench FinOps task; analyze local CSV/table files only.",
        "Do not use Kubernetes SRE tools for FinOps scenarios.",
        "Use analyze_finops_cost_anomaly first; it normalizes date formats before matching.",
        "Return only resource identifiers in FINAL_ANSWER, comma-separated.",
    ],
    "itbench_lite": [
        "Use only ITBench SRE repository-aligned tools in this ITBench-Lite mode.",
        "Do not call generic web tools in ITBench mode.",
        "Start with summarize_sre_candidates, then use rank_sre_root_cause_candidates before final answer.",
        "You may use read_file for deterministic parsing of local snapshot files.",
        "Do not answer with namespaces; answer with entity IDs or entity names.",
        "Return only resource/group identifiers in FINAL_ANSWER, comma-separated.",
    ],
    "itbench_lite_sre": [
        "This is an ITBench SRE task; use only SRE repository-aligned tools.",
        "Do not use FinOps CSV/Python tools for SRE scenarios.",
        "Start with summarize_sre_candidates, then use rank_sre_root_cause_candidates before final answer.",
        "If ranking is ambiguous, use read_file on k8s_events_raw.tsv and k8s_objects_raw.tsv to verify entity names.",
        "Do not answer with namespaces; answer with entity IDs or entity names.",
        "Return only resource/group identifiers in FINAL_ANSWER, comma-separated.",
    ],
}


def normalize_tool_name(name: str) -> str:
    n = (name or "").strip().lower()
    n = n.replace("-", "_").replace(" ", "_").replace("/", "_")
    n = __import__("re").sub(r"[^a-z0-9_]+", "", n)
    alias = {
        "pdf_viewer": "read_file",
        "pdf_access": "read_file",
        "excel": "read_file",
        "microsoft_excel": "read_file",
        "excel_file_access": "read_file",
        "word_document_access": "read_file",
        "xml_file_access": "read_file",
        "jsonld_file_access": "read_file",
        "image_recognition": "read_file",
        "image_recognition_tools": "read_file",
        "ocr": "read_file",
        "tool_to_extract_text_from_images": "read_file",
        "file_handling": "read_file",
        "pdf_reader": "read_file",
        "pdf_readerextracter": "read_file",
        "xlsx_file_access": "read_file",
        "csv_file_access": "read_file",
        "access_to_excel_files": "read_file",
        "spreadsheet_editor": "read_file",
        "powerpoint_viewer": "read_file",
        "color_recognition": "read_file",
        "csv_summary": "summarize_csv",
        "summarize_csv": "summarize_csv",
        "nl2_kubectl": "nl2kubectl",
        "kubectl": "nl2kubectl",
        "loki_query": "query_loki_logs",
        "jaeger_query": "query_jaeger_traces",
        "topology_nodes": "get_topology_nodes",
        "sre_candidate_summary": "summarize_sre_candidates",
        "summarize_candidates": "summarize_sre_candidates",
        "root_cause_candidate_rank": "rank_sre_root_cause_candidates",
        "rank_root_causes": "rank_sre_root_cause_candidates",
        "finops_anomaly": "analyze_finops_cost_anomaly",
        "analyze_finops_anomaly": "analyze_finops_cost_anomaly",
        "finops_cost_anomaly": "analyze_finops_cost_anomaly",
        "node_info_by_name": "get_node_info_by_name",
    }
    if n in alias:
        return alias[n]
    if "excel" in n or "spreadsheet" in n or "csv" in n or "xlsx" in n:
        return "read_file"
    if "pdf" in n or "word" in n or "xml" in n or "json" in n or "file" in n:
        return "read_file"
    if "image" in n or "ocr" in n or "vision" in n or "gif" in n:
        return "read_file"
    if "kubectl" in n or "k8s" in n or "kubernetes" in n:
        return "nl2kubectl"
    if "loki" in n:
        return "query_loki_logs"
    if "jaeger" in n or "trace" in n:
        return "query_jaeger_traces"
    if "alert" in n:
        return "get_alerts"
    if "topology" in n:
        return "get_topology_nodes"
    if "neighbor" in n:
        return "get_neighbors"
    if "connected" in n or "connectivity" in n:
        return "check_directly_connected"
    if "finops" in n or "anomaly" in n or "cost" in n:
        return "analyze_finops_cost_anomaly"
    return n


def _parse_tool_prompt_profile(profile: str) -> tuple[str, str]:
    raw = str(profile or "native")
    if ":" in raw:
        base, bench = raw.split(":", 1)
    else:
        base, bench = raw, "itbench_lite"
    base = base.strip().lower() or "native"
    if base not in {"native", "simple"}:
        base = "native"
    bench = bench.strip().lower() or "itbench_lite"
    if bench not in BENCHMARK_TOOLSETS:
        bench = "itbench_lite"
    return base, bench


def toolset_for_profile(profile: str) -> list[str]:
    _, bench = _parse_tool_prompt_profile(profile)
    return sorted(BENCHMARK_TOOLSETS.get(bench, BENCHMARK_TOOLSETS["itbench_lite"]))


BASE_TOOL_PROMPT = """You are an agent that can call tools.
You must respond in exactly one of these formats:
1) TOOL_CALL: {"tool":"tool_name","args":{...}}
2) FINAL_ANSWER: <answer>

Rules:
- One tool call per turn.
- No markdown, no code fences, no explanation.
- JSON must be valid and minimal.
- Available tools are listed in the schema section below; use only those tools.
"""


SIMPLE_TOOL_PROMPT = """You are a strict tool-calling assistant.
Return exactly one line in one of these forms:
1) TOOL_CALL: {"tool":"tool_name","args":{...}}
2) FINAL_ANSWER: <answer>

Rules:
- One tool call per turn.
- No markdown, no code fences, no explanation.
- JSON must be valid and minimal.
- Available tools are listed in the schema section below; use only those tools.
"""


def build_tool_prompt(profile: str = "native") -> str:
    base_profile, bench = _parse_tool_prompt_profile(profile)
    allowed = set(toolset_for_profile(profile))
    base = SIMPLE_TOOL_PROMPT if base_profile == "simple" else BASE_TOOL_PROMPT
    lines = [base.strip(), "", "Tool schemas (expected args):"]
    for name, schema in TOOL_SCHEMAS.items():
        if name not in allowed:
            continue
        desc = TOOL_DESCRIPTIONS.get(name, "")
        lines.append(
            f"- {name}: {json.dumps(schema, ensure_ascii=False)}"
            + (f" | {desc}" if desc else "")
        )
    lines.append("")
    lines.append(
        "If uncertain about argument keys, still call the tool with your best guess; argument aliases are supported."
    )
    lines.append(
        "Important: when a schema says Python syntax, provide Python-style expressions/code exactly."
    )
    guidance = BENCHMARK_TOOL_GUIDANCE.get(bench, [])
    if guidance:
        lines.append("")
        lines.append("Benchmark-specific guidance:")
        for g in guidance:
            lines.append(f"- {g}")
    return "\n".join(lines)


def tool_schema_hint(tool: str) -> str:
    t = normalize_tool_name(tool or "")
    if t in TOOL_SCHEMAS:
        return f"Expected {t} args schema: {json.dumps(TOOL_SCHEMAS[t], ensure_ascii=False)}"
    available = ", ".join(sorted(TOOL_SCHEMAS.keys()))
    return f"Unknown tool '{tool}'. Available tools: {available}"


def _adapter(tool_name: str):
    def _call(**kwargs) -> str:
        return run_tool(tool_name, kwargs or {})

    return _call


TOOLS = {
    "read_file": read_file,
    "summarize_csv": summarize_csv,
    "nl2kubectl": _adapter("nl2kubectl"),
    "query_loki_logs": _adapter("query_loki_logs"),
    "query_jaeger_traces": _adapter("query_jaeger_traces"),
    "get_alerts": _adapter("get_alerts"),
    "get_topology_nodes": _adapter("get_topology_nodes"),
    "walk_path": _adapter("walk_path"),
    "get_node_info_by_name": _adapter("get_node_info_by_name"),
    "get_neighbors": _adapter("get_neighbors"),
    "check_directly_connected": _adapter("check_directly_connected"),
    "analyze_finops_cost_anomaly": _adapter("analyze_finops_cost_anomaly"),
    "summarize_sre_candidates": _adapter("summarize_sre_candidates"),
    "rank_sre_root_cause_candidates": _adapter("rank_sre_root_cause_candidates"),
}
